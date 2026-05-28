import json
import sys
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation

from docgraph_sidecar.core.db import connect, initialize_database
from docgraph_sidecar.parser import (
    ParseContext,
    ParserError,
    PptxParser,
    default_parser_registry,
    parse_with_error_recording,
)


class PptxParserTest(unittest.TestCase):
    def setUp(self) -> None:
        self.temp_dir = TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        self.parser = PptxParser()

    def tearDown(self) -> None:
        self.temp_dir.cleanup()

    def test_pptx_parser_extracts_title_body_and_notes_per_slide(self) -> None:
        path = self.root / "status.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Alpha Status"
        slide.placeholders[1].text = "Roadmap is on track.\nDelivery notes are ready."
        slide.notes_slide.notes_text_frame.text = "Discuss camera rollout."
        presentation.save(path)

        result = self.parser.parse(ParseContext.from_path(path, file_id="file-pptx"))

        self.assertEqual(result.parser_name, "pptx")
        self.assertEqual(len(result.elements), 1)
        element = result.elements[0]
        self.assertEqual(element.element_type, "slide")
        self.assertEqual(element.slide_no, 1)
        self.assertEqual(element.metadata["title"], "Alpha Status")
        self.assertEqual(element.metadata["body"], ["Roadmap is on track. Delivery notes are ready."])
        self.assertEqual(element.metadata["notes"], "Discuss camera rollout.")
        self.assertFalse(element.metadata["ocr_performed"])
        self.assertIn("Title: Alpha Status", result.chunks[0].text)
        self.assertIn("Notes: Discuss camera rollout.", result.chunks[0].text)
        self.assertEqual(result.chunks[0].heading, "Alpha Status")
        self.assertFalse(result.chunks[0].evidence["ocr_performed"])

    def test_multiple_slides_create_multiple_chunks(self) -> None:
        path = self.root / "multi.pptx"
        presentation = Presentation()
        first = presentation.slides.add_slide(presentation.slide_layouts[1])
        first.shapes.title.text = "Summary"
        first.placeholders[1].text = "Alpha project summary"
        second = presentation.slides.add_slide(presentation.slide_layouts[1])
        second.shapes.title.text = "Risks"
        second.placeholders[1].text = "No image OCR in this milestone"
        presentation.save(path)

        result = self.parser.parse(ParseContext.from_path(path, file_id="file-multi"))

        self.assertEqual([chunk.slide_no for chunk in result.chunks], [1, 2])
        self.assertEqual([chunk.heading for chunk in result.chunks], ["Summary", "Risks"])
        self.assertEqual(result.metadata["presentation_profile"]["slide_count"], 2)

    def test_basic_fixture_pptx_can_be_parsed(self) -> None:
        path = Path("fixtures/basic_docs/alpha_status.pptx")

        result = self.parser.parse(ParseContext.from_path(path, file_id="file-fixture"))

        self.assertEqual(len(result.elements), 1)
        self.assertEqual(result.elements[0].slide_no, 1)
        self.assertIn("Alpha Project status deck", result.chunks[0].text)

    def test_default_registry_selects_pptx_parser(self) -> None:
        path = self.root / "simple.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Simple deck"
        slide.placeholders[1].text = "Body"
        presentation.save(path)

        result = default_parser_registry().parse_path(path, file_id="file-registry")

        self.assertEqual(result.parser_name, "pptx")
        self.assertEqual(result.chunks[0].heading, "Simple deck")

    def test_bad_pptx_records_parse_error(self) -> None:
        data_dir = self.root / "data"
        bad_path = self.root / "bad.pptx"
        bad_path.write_bytes(b"not a presentation")
        initialize_database(data_dir=data_dir)
        self._insert_file(data_dir, "file-bad", str(bad_path))
        context = ParseContext.from_path(bad_path, file_id="file-bad")

        with self.assertRaises(ParserError):
            parse_with_error_recording(default_parser_registry(), context, data_dir=data_dir)

        connection = connect(data_dir=data_dir)
        try:
            row = connection.execute(
                """
                SELECT error_code, error_message, parser_name, retryable, details_json
                FROM parse_errors
                WHERE file_id = 'file-bad'
                """
            ).fetchone()
        finally:
            connection.close()

        self.assertIsNotNone(row)
        self.assertEqual(row["error_code"], "PPTX_PARSE_ERROR")
        self.assertEqual(row["parser_name"], "pptx")
        self.assertEqual(row["retryable"], 0)
        self.assertEqual(row["error_message"], "PPTX file could not be parsed.")
        self.assertEqual(json.loads(row["details_json"])["path"], str(bad_path))

    def _insert_file(self, data_dir: Path, file_id: str, path: str) -> None:
        connection = connect(data_dir=data_dir)
        try:
            connection.execute(
                """
                INSERT INTO files (
                  file_id,
                  path,
                  normalized_path,
                  filename,
                  extension,
                  source_type,
                  file_status,
                  parse_status,
                  deleted_flag
                )
                VALUES (?, ?, ?, ?, '.pptx', 'office', 'discovered', 'pending', 0)
                """,
                (file_id, path, path.casefold(), Path(path).name),
            )
            connection.commit()
        finally:
            connection.close()


if __name__ == "__main__":
    unittest.main()
