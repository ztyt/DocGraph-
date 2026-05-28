from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docgraph_sidecar.parser.base import (
    BaseParser,
    ParsedChunk,
    ParsedDocumentElement,
    ParseContext,
    ParseResult,
    ParserError,
)


@dataclass(frozen=True)
class SlideProfile:
    slide_no: int
    title: str | None
    body: tuple[str, ...]
    notes: str | None
    image_count: int

    def to_metadata(self) -> dict[str, Any]:
        return {
            "slide_no": self.slide_no,
            "title": self.title,
            "body": list(self.body),
            "notes": self.notes,
            "image_count": self.image_count,
            "ocr_performed": False,
        }


class PptxParser(BaseParser):
    name = "pptx"
    supported_extensions = (".pptx",)

    def parse(self, context: ParseContext) -> ParseResult:
        try:
            presentation = Presentation(context.path)
        except Exception as exc:
            raise ParserError(
                "PPTX file could not be parsed.",
                error_code="PPTX_PARSE_ERROR",
                parser_name=self.name,
                retryable=False,
                details={"path": str(context.path), "error_type": type(exc).__name__},
            ) from exc

        elements: list[ParsedDocumentElement] = []
        chunks: list[ParsedChunk] = []
        profiles: list[SlideProfile] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            profile = profile_slide(slide, slide_no=slide_index)
            profiles.append(profile)
            text = slide_profile_text(profile)
            if not text:
                continue
            add_slide_block(elements, chunks, context=context, profile=profile, text=text)

        warnings = () if elements else ("No parseable PPTX slides found.",)
        return ParseResult(
            parser_name=self.name,
            file_id=context.file_id,
            elements=tuple(elements),
            chunks=tuple(chunks),
            warnings=warnings,
            metadata={
                "presentation_profile": {
                    "slide_count": len(profiles),
                    "slides": [profile.to_metadata() for profile in profiles],
                }
            },
        )


def profile_slide(slide: Any, *, slide_no: int) -> SlideProfile:
    title_shape = slide.shapes.title
    title = normalize_text(title_shape.text) if title_shape is not None else None
    body: list[str] = []
    image_count = 0

    for shape in slide.shapes:
        if is_picture_shape(shape):
            image_count += 1
        if title_shape is not None and shape == title_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = normalize_text(shape.text)
        if text:
            body.append(text)

    notes = None
    if getattr(slide, "has_notes_slide", False):
        notes_text = normalize_text(slide.notes_slide.notes_text_frame.text)
        notes = notes_text or None

    return SlideProfile(
        slide_no=slide_no,
        title=title,
        body=tuple(body),
        notes=notes,
        image_count=image_count,
    )


def slide_profile_text(profile: SlideProfile) -> str:
    lines = [f"Slide {profile.slide_no}"]
    if profile.title:
        lines.append(f"Title: {profile.title}")
    if profile.body:
        lines.append("Body:")
        lines.extend(profile.body)
    if profile.notes:
        lines.append(f"Notes: {profile.notes}")
    return "\n".join(lines).strip()


def is_picture_shape(shape: Any) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def add_slide_block(
    elements: list[ParsedDocumentElement],
    chunks: list[ParsedChunk],
    *,
    context: ParseContext,
    profile: SlideProfile,
    text: str,
) -> None:
    index = len(elements)
    element_id = stable_parse_id(context.file_id, "element", index, text)
    chunk_id = stable_parse_id(context.file_id, "chunk", index, text)
    metadata = profile.to_metadata()
    elements.append(
        ParsedDocumentElement(
            element_id=element_id,
            file_id=context.file_id,
            element_index=index,
            element_type="slide",
            slide_no=profile.slide_no,
            section_path=profile.title,
            text=text,
            metadata=metadata,
            confidence=1.0,
        )
    )
    chunks.append(
        ParsedChunk(
            chunk_id=chunk_id,
            file_id=context.file_id,
            element_id=element_id,
            chunk_index=index,
            chunk_type="slide",
            slide_no=profile.slide_no,
            heading=profile.title,
            section_path=profile.title,
            text=text,
            token_count=estimate_token_count(text),
            evidence={"parser": PptxParser.name, "ocr_performed": False},
        )
    )


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", value.replace("\xa0", " ")).strip()


def stable_parse_id(file_id: str, kind: str, index: int, text: str) -> str:
    digest = hashlib.sha256(f"{file_id}:{kind}:{index}:{text}".encode("utf-8")).hexdigest()
    return f"{kind}-{digest[:24]}"


def estimate_token_count(text: str) -> int:
    return len(re.findall(r"\S+", text))
